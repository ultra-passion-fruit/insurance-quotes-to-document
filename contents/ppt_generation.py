from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from datetime import datetime
from pandas import read_csv
from quotes_calculation import get_access_token, get_quote
import os

# Date: 202509009
# Modified: 20260113

# This script will generate a PowerPoint (pptx) presentation for Energia Insurance.
# The goal is to generate a PowerPoint presentation with insurance quote information from the insurance company.
# The insurance company quote information is to be accessed with the insurance company API.

FONT = 'Segoe UI'

def print_shapes(slide):
    """
    Prints all shapes on a slides
    
    :param slide: Description
    """
    index = 0
    for shape in slide.shapes:
        print("\nIndex: ",index)
        print("Has Table: ",shape.has_table)
        print("Has Text Frame: ", shape.has_text_frame)
        if shape.has_text_frame:
            print("Text Value: ",shape.text_frame.text)
        print("Object: ",shape)
        index = index + 1

def set_text(run, text, font_name, size, color='Black', bold=False, italic=False):
    """
    Sets the text of a 'run' with style parameters
    Returns: a run
    """
    run.text = text
    # setting font style
    font = run.font
    font.name = font_name
    font.size = size
    font.color.theme_color = color
    font.bold = bold
    font.italic = italic
    return run

def set_cover(text_frame, text):
    """
    Sets the text contents of a text frame object
    
    :param text_frame: text frame object
    :param text: text to write
    """
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    set_text(run, text, FONT, Pt(11), color=MSO_THEME_COLOR.ACCENT_6)

def fill_table(table,some_list):
    for x in range(0,len(some_list)):
        # get cell in the second row, first column (i.e., skip table header)
        coverage_textframe = table.cell(x+1,0).text_frame
        # clear whatever is in text frame
        coverage_textframe.clear()
        # by default always have 1 paragraph when empty, grab that paragraph
        paragraph = coverage_textframe.paragraphs[0]
        # add a new run to it
        run = paragraph.add_run()
        # set the text to be added
        coverage_text = f"${some_list[x][0]:,}"
        # setting text value
        set_text(run, coverage_text, 'Segoe UI', Pt(10.5), color=MSO_THEME_COLOR.TEXT_1)
        
        # get cell in second column, the second row (i.e., skip table header)
        quote_textframe = table.cell(x+1,1).text_frame
        # clear whatever is in text frame
        quote_textframe.clear()
        # by default always have 1 paragraph when empty, grab that paragraph
        paragraph = quote_textframe.paragraphs[0]
        # add a new run to it
        run = paragraph.add_run()
        # set the text to be added
        quote_text = "$" + str(some_list[x][1])
        # setting text value
        set_text(run, quote_text, 'Segoe UI', Pt(10.5), color=MSO_THEME_COLOR.TEXT_1)

def print_table(table):
    for row in range(len(table.rows)):
        for col in range(len(table.columns)):
            print(table.cell(row, col).text, end=" ")
        print()

###################################################################################################

##### Getting data from .csv #####

# change into contents (when running in terminal runs from parent)
os.chdir('contents')

# df from jotform
df = read_csv('client-jotform.csv')
# select relevant columns to keep from jotform
columns_to_keep = [df.columns[1], df.columns[2], df.columns[5], df.columns[34], df.columns[35], df.columns[36], df.columns[29], df.columns[30], df.columns[31], df.columns[6], df.columns[7], df.columns[8], df.columns[12], df.columns[13], df.columns[16], df.columns[17], df.columns[20], df.columns[21], df.columns[24], df.columns[25]]
client_df = df[columns_to_keep]
# renaming columns
client_df.columns = ["full_name", "email", "submit", "trip_start", "dest_arrival", "end", "dest_country", "dest_region", "depart_region", "num_travellers", "name_1", "dob_1", "name_2", "dob_2", "name_3", "dob_3", "name_4", "dob_4", "name_5", "dob_5"]
client_dict = client_df.to_dict('records')[0]

# Reformatting dates
old_format = "%b %d, %Y"
new_format = "%Y-%m-%d"

# Reformatting start and end dates
submit_string = client_dict["submit"] # get string
submit_date_obj = datetime.strptime(submit_string, old_format) #create obj
depart_string = client_dict["trip_start"] # get string
depart_date_obj = datetime.strptime(depart_string, old_format) #create obj
start_string = client_dict["dest_arrival"] # get string
start_date_obj = datetime.strptime(start_string, old_format) #create obj
end_string = client_dict["end"] # get string
end_date_obj = datetime.strptime(end_string, old_format)  #created obj
# setting new format
client_dict["submit"] = submit_date_obj.strftime(new_format)
client_dict["trip_start"] = depart_date_obj.strftime(new_format)
client_dict["dest_arrival"] = start_date_obj.strftime(new_format)
client_dict["end"] = end_date_obj.strftime(new_format)

# Reformatting dob dates
for x in range(1,client_dict["num_travellers"]+1):
    dob_string = client_dict[f"dob_{x}"]
    dob_date_obj = datetime.strptime(dob_string, old_format) #create obj
    client_dict[f"dob_{x}"] = dob_date_obj.strftime(new_format) # setting new format

###################################

### Cover Formatting Info ###
# Font Size: 11
# Font Name: Segoe UI
# Bold: Yes
# Font Color: 172, 110, 223

##### Slide manipulation / formatting #######

prs = Presentation('template.pptx')
purple = RGBColor.from_string("ac6edf")

print("\nAdding client info to cover page...")

# get first slide
cover_slide = prs.slides[0]

# setting client name on cover
name_text_frame = cover_slide.shapes[1].text_frame
set_cover(name_text_frame, client_dict["full_name"])

# setting client date of birth on cover
dob_text_frame = cover_slide.shapes[2].text_frame
set_cover(dob_text_frame, client_dict["dob_1"])

# setting client email on cover
email_text_frame = cover_slide.shapes[3].text_frame
set_cover(email_text_frame, client_dict["email"])

# setting client destination on cover
destination_text_frame = cover_slide.shapes[4].text_frame
set_cover(destination_text_frame, client_dict["dest_country"])

print("Client information added to cover page:")
for i in range(1,5):
    print(prs.slides[0].shapes[i].text_frame.text)

############ Requesting Quote Data from API and formatting slide's table #############

print("\nGetting quotes information from Tugo...")

# get slide
tugo_slide = prs.slides[2]

# get the table from the slide
tugo_quotes = tugo_slide.shapes[6].table

# get access token to request quotes price
access_token = get_access_token()

# list of all coverage levels
coverage_levels = [10000, 25000, 50000, 100000, 200000, 300000, 500000]

# Create list with quotes for each coverage levels
    # inner list represents a row in table
    # value in each inner list is value in each row
quotes_list = []
for coverage in coverage_levels:
    try:
        quote = get_quote(coverage, access_token, client_dict)['availablePlanPrices'][0]['ratedPrice']['total']
        table_row = [coverage, quote]
        quotes_list.append(table_row)
    except Exception as e:
        raise e

# filling table
fill_table(tugo_quotes,quotes_list)

# printing to console
print("Quotes information added to TuGo slide:")
print_table(tugo_quotes)

# Setting info on Tugo slide
tugo_info_header = tugo_slide.shapes[25].table

# setting seguro dates
tugo_dates_text_frame = tugo_info_header.cell(0,1).text_frame
tugo_dates_text_frame.clear()
d_para = tugo_dates_text_frame.paragraphs[0]
d_run = d_para.add_run()
d_text = client_dict["trip_start"] + " - " + client_dict["end"]
set_text(d_run, d_text, FONT, Pt(9), MSO_THEME_COLOR.ACCENT_2)

print("\nAdding header information to TuGo slide:")

# setting number of segurados
tugo_insured_text_frame = tugo_info_header.cell(0,3).text_frame
tugo_insured_text_frame.clear()
d_para = tugo_insured_text_frame.paragraphs[0]
d_run = d_para.add_run()
num_travelers = client_dict["num_travellers"]
d_text = str(num_travelers) + " pessoa" + 's'*(num_travelers-1)
set_text(d_run, d_text, FONT, Pt(9), color=MSO_THEME_COLOR.ACCENT_2)

# printing to console
print("Tugo Header Information Added:")
print_table(tugo_info_header)

##################################################################################

# Saving Presentation

savedir = 'quoted_presentations'
savefile = f'{client_dict["full_name"].replace(' ', '-')}.pptx'

print("\nSaving presentation...")
print(f"Saved to '{savefile}' in {savedir} folder.")

os.chdir('..')

if savedir not in os.listdir():
    os.mkdir(savedir)

os.chdir('quoted_presentations')

# saving as pptx
prs.save(savefile)
# prs.save(f'../quoted_pptx/{client_dict["full_name"]}.pptx')


