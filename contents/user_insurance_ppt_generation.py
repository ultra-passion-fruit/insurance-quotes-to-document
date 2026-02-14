# Author: Matsuru Hoshi
# Date: 202509009

# This script will generate a PowerPoint (pptx) presentation for Energia Insurance.
# The goal is to generate a PowerPoint presentation with insurance quote information from the insurance company.
# The insurance company quote information is to be accessed with the companies' respectable APIs.

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor



person = {"name": "Paulo Silva",
          "dob": "1900-01-01",
          "email": "paulo@gmail.com",
          "destination": "Canadá"}
trip_info = {"insurance_start": "2025-01-01",
          "insurance_end": "2025-02-1"}

company = {"name": "Company A",
           "policy_link": "https://energiaseguros.com/",
           "policy_price": "100.99"}

prs = Presentation('template.pptx')
purple = RGBColor.from_string("ac6edf")
# cover_info = prs.slides[0].shapes[1].table
# cover_info_cell = prs.slides[0].shapes[1].table.cell(0,0).text

### Cover Formatting Info ###
# Font Size: 11
# Font Name: Segoe UI
# Bold: Yes
# Font Color: 172, 110, 223

p1 = prs.slides[0].shapes[3]
# print(p1.text)

i = 1
for key in person:
    p1 = prs.slides[0].shapes[i]
    p1.text = ""
    p1_tf = p1.text_frame.add_paragraph()
    p1_tf.text = person[key]
    p1_tf.font.size = Pt(11) 
    p1_tf.font.color.rgb = purple
    i = i + 1

    print(p1.text)

tugo_slide = prs.slides[2]

def print_shapes(slide):
    index = 0
    for shape in slide.shapes:
        print("Index: {ind}\nHas Table:{tab}\nHas Text:{tex}\nObject:{obj}\n".format(ind=index, tab=shape.has_table, tex=shape.has_text_frame, obj=shape))
        index = index + 1


tugo_quotes = tugo_slide.shapes[6].table

print(tugo_quotes.cell(0,0).text)

prs.save('template.pptx')
